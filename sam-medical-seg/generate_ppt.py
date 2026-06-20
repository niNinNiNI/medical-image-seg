#!/usr/bin/env python3
"""
Generate the SAM-MedSeg presentation slides (PPTX).
Covers: title, motivation, method, experiments, results, conclusion.

Usage:
    python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent
PAPER_DIR = PROJECT_ROOT / "paper"
PAPER_DIR.mkdir(parents=True, exist_ok=True)
FIGURES_DIR = PROJECT_ROOT / "experiments" / "figures"

# Color palette
BLUE = RGBColor(0x2C, 0x7B, 0xB6)
RED = RGBColor(0xD7, 0x19, 0x1C)
GREEN = RGBColor(0x1A, 0x96, 0x41)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)


def add_slide_title(prs, title, subtitle=None):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.0))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, notes=None):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Divider line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(1.05), Inches(9), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # Bullet points
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.level = 0

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_image_slide(prs, title, image_path, width_inches=8):
    """Add a slide with an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches((10 - width_inches) / 2),
            Inches(1.0),
            Inches(width_inches),
        )
    return slide


def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══ Slide 1: Title ═══
    add_slide_title(prs,
        "SAM-MedSeg: Improving SAM for\nFew-Shot Medical Image Segmentation",
        "Deep Learning & Computer Vision — Course Project\nJune 2026")

    # ═══ Slide 2: Motivation ═══
    add_content_slide(prs, "Motivation: The Annotation Bottleneck", [
        "• Medical image segmentation requires pixel-level annotations",
        "• Annotation needs radiologist expertise → expensive & time-consuming",
        "• Few-shot segmentation: achieve good performance with very limited labels",
        "• SAM shows impressive zero-shot on natural images — can we leverage it?",
        "",
        "⚠ SAM zero-shot on Kvasir-SEG polyps:",
        "   Dice = 0.271, Precision = 0.175, Recall = 0.921",
        "   → High recall but extreme false positives — clinically unusable",
    ])

    # ═══ Slide 3: Proposed Method Overview ═══
    add_content_slide(prs, "SAM-MedSeg: Three Key Components", [
        "① LoRA Adaptation (227K params, 0.25% of backbone)",
        "   • Low-rank decomposition applied to SAM's QKV/proj matrices",
        "   • Rank r=4 — injects domain knowledge with minimal overhead",
        "",
        "② Automatic Prompt Optimization",
        "   • 32×32 grid sampling → per-point SAM inference → IoU filtering → NMS",
        "   • Fully automatic, no manual prompts needed",
        "",
        "③ Medical Skip-Connection Decoder (30.4M params)",
        "   • U-Net style decoder with skip connections from SAM encoder layers",
        "   • Recovers fine boundary details lost in SAM's single-scale decoder",
        "",
        "④ Learnable Fusion: M_final = σ(α)·M_sam + (1−σ(α))·M_med  (α=0.623)",
    ])

    # ═══ Slide 4: Architecture Figure ═══
    arch_fig = FIGURES_DIR / "figure_arch_1_overall.png"
    add_image_slide(prs, "Overall Architecture", arch_fig, 8.5)

    # ═══ Slide 5: Experimental Setup ═══
    add_content_slide(prs, "Experimental Setup", [
        "📊 Dataset: Kvasir-SEG (1,000 polyp images, 1024×1024)",
        "   Train/Val/Test = 800/100/100",
        "",
        "🔬 Data Efficiency Settings:",
        "   1% (8 imgs) | 5% (40 imgs) | 10% (80 imgs) | 100% (800 imgs)",
        "   3 seeds each (42, 43, 44) → 12 + 12 + 1 + 2 = 27 experiments total",
        "",
        "📏 Metrics: Dice, IoU, Precision, Recall",
        "",
        "🖥 Hardware: Single NVIDIA RTX 4060 Laptop (8GB VRAM)",
        "   AMP fp16 + Gradient Checkpointing → Peak 6.16 GB",
        "",
        "🔬 Baselines: SAM Zero-Shot | U-Net from scratch | LoRA Only | MedDecoder Only",
    ])

    # ═══ Slide 6: Main Results ═══
    data_fig = FIGURES_DIR / "figure_1_data_efficiency.png"
    add_image_slide(prs, "Data Efficiency: SAM-MedSeg vs U-Net", data_fig, 8.5)

    # ═══ Slide 7: Crossover Analysis ═══
    add_content_slide(prs, "Key Finding: The Crossover Point", [
        "🎯 Performance crossover between 5% and 10% data",
        "",
        "Data % | U-Net | SAM-MedSeg | Winner      | Gap",
        "───────┼───────┼────────────┼─────────────┼──────",
        "  1%   | 0.385 |   0.425    | SAM-MedSeg ✅ | +4.0 pp",
        "  5%   | 0.443 |   0.467    | SAM-MedSeg ✅ | +2.4 pp",
        " 10%   | 0.509 |   0.484    | U-Net 🔄     | +2.5 pp",
        " 100%  | 0.719 |   0.640    | U-Net        | +7.9 pp",
        "",
        "💡 Practical takeaway:",
        "  ≤50 annotated images → Use SAM-MedSeg with LoRA",
        "  >50 annotated images → U-Net from scratch is better & simpler",
    ])

    # ═══ Slide 8: Ablation ═══
    ablation_fig = FIGURES_DIR / "figure_2_ablation_study.png"
    add_image_slide(prs, "Ablation Study: Component Contributions", ablation_fig, 8.5)

    add_content_slide(prs, "Ablation Insights", [
        "🔬 At 5% data (40 images, seed 42):",
        "",
        "   Configuration      | Dice  | Precision | Recall | Params",
        "   ───────────────────┼───────┼───────────┼────────┼───────",
        "   Full Model          | 0.448 |   0.421   | 0.697  | 4.84M",
        "   LoRA Only           | 0.452 |   0.441   | 0.658  | 227K  👑",
        "   MedDecoder Only     | 0.426 |   0.366   | 0.772  | 30.4M",
        "",
        "💡 LoRA is the core driver of few-shot generalization",
        "   • 227K params alone achieve best Dice (0.452)",
        "   • MedDecoder provides complementary high recall (0.772)",
        "   • Full model balances both strengths via learnable fusion",
    ])

    # ═══ Slide 9: Segmentation Examples ═══
    seg_fig = FIGURES_DIR / "figure_3_segmentation_comparison.png"
    add_image_slide(prs, "Visual Comparison: Segmentation Examples", seg_fig, 8.5)

    # ═══ Slide 10: Training Curves ═══
    train_fig = FIGURES_DIR / "figure_4_training_curves.png"
    add_image_slide(prs, "Training Curves: Best Models", train_fig, 8.5)

    # ═══ Slide 11: Limitations ═══
    add_content_slide(prs, "Limitations & Future Work", [
        "⚠ Known Limitations:",
        "   • Over-segmentation of tiny polyps (<3% foreground, Dice 0.237)",
        "   • Only validated on single 2D dataset (Kvasir-SEG)",
        "   • Grid-sampling prompt strategy is computationally expensive (1024 forward passes)",
        "   • Fixed LoRA rank (r=4) — not optimized per data scale",
        "",
        "🔮 Future Directions:",
        "   • Extend to 3D medical volumes (CT/MRI) with 3D SAM + 3D LoRA",
        "   • Lightweight prompt proposal network to replace exhaustive grid search",
        "   • Active learning for optimal sample selection within annotation budget",
        "   • Multi-center, multi-modal validation for robustness assessment",
        "   • Incorporate clinical metadata via text encoders for context-aware segmentation",
    ])

    # ═══ Slide 12: Conclusion ═══
    add_content_slide(prs, "Conclusion", [
        "✅ SAM-MedSeg: parameter-efficient framework for few-shot medical segmentation",
        "",
        "✅ LoRA (227K params) is the dominant driver — not expensive decoders",
        "",
        "✅ SAM-MedSeg is OPTIMAL when labeled data ≤5% (≤50 images)",
        "",
        "✅ Crossover at 5-10% — above this, traditional U-Net is better",
        "",
        "✅ Provides practical model selection guidance based on annotation budget",
        "",
        "✅ All code, experiments, and figures are reproducible & publicly available",
        "",
        "📄 Paper: ~4,400 words | 📊 8 figures | 📚 17 references | 🔬 27 experiments",
    ])

    # ═══ Slide 13: Thank You ═══
    add_slide_title(prs, "Thank You!\nQuestions?",
        "Code & Paper: [GitHub Repository]\nContact: [Team Members]")

    # Save
    output = PAPER_DIR / "SAM-MedSeg_presentation.pptx"
    prs.save(str(output))
    print(f"PPT saved to: {output}")
    return output


if __name__ == "__main__":
    generate_ppt()
